#!/usr/bin/env python3
"""
Script to apply Infoblox styling to PowerPoint presentation
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# Define Infoblox brand colors
# Primary colors
INFOBLOX_BLUE = RGBColor(0, 125, 186)  # #007DBA - Primary blue
INFOBLOX_DARK_BLUE = RGBColor(0, 63, 94)  # #003F5E - Dark blue
INFOBLOX_LIGHT_BLUE = RGBColor(0, 176, 240)  # #00B0F0 - Light blue
INFOBLOX_GRAY = RGBColor(88, 88, 91)  # #58585B - Dark gray
INFOBLOX_LIGHT_GRAY = RGBColor(236, 236, 236)  # #ECECEC - Light gray background

# Load the existing presentation
presentation_path = '/home/parallels/Work/Tutorials/Hackathon/2025/Kea_DHCP_Monitoring_Presentation.pptx'
prs = Presentation(presentation_path)

# Create a new presentation with Infoblox styling
infoblox_prs = Presentation()
infoblox_prs.slide_width = Inches(13.33)
infoblox_prs.slide_height = Inches(7.5)

# Function to add an Infoblox footer to a slide
def add_infoblox_footer(slide):
    # Add a thin line above the footer
    line_shape = slide.shapes.add_shape(
        1, Inches(0.5), Inches(6.9), Inches(12.33), Inches(0.02)
    )
    line_shape.fill.solid()
    line_shape.fill.fore_color.rgb = INFOBLOX_BLUE
    line_shape.line.fill.background()

    # Add footer text
    footer = slide.shapes.add_textbox(Inches(0.5), Inches(7.0), Inches(8), Inches(0.3))
    footer_text = footer.text_frame
    footer_text.paragraphs[0].text = "Kea DHCP Server Monitoring | May 10, 2025"
    footer_text.paragraphs[0].font.size = Pt(10)
    footer_text.paragraphs[0].font.color.rgb = INFOBLOX_GRAY

    # Add logo placeholder (right side of footer)
    logo_placeholder = slide.shapes.add_textbox(Inches(11), Inches(7.0), Inches(1.8), Inches(0.3))
    logo_text = logo_placeholder.text_frame
    logo_text.paragraphs[0].text = "INFOBLOX"
    logo_text.paragraphs[0].alignment = PP_ALIGN.RIGHT
    logo_text.paragraphs[0].font.bold = True
    logo_text.paragraphs[0].font.size = Pt(12)
    logo_text.paragraphs[0].font.color.rgb = INFOBLOX_BLUE

# Function to create Infoblox styled title slide
def create_title_slide(content_slide):
    slide_layout = infoblox_prs.slide_layouts[0]  # Title Slide layout
    slide = infoblox_prs.slides.add_slide(slide_layout)

    # Add blue background shape at the top
    top_bg = slide.shapes.add_shape(
        1, Inches(0), Inches(0), Inches(13.33), Inches(2.5)
    )
    top_bg.fill.solid()
    top_bg.fill.fore_color.rgb = INFOBLOX_BLUE
    top_bg.line.fill.background()

    # Title text
    title_box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(11.33), Inches(1))
    title_text = title_box.text_frame
    title_text.word_wrap = True

    p = title_text.paragraphs[0]
    p.text = "Kea DHCP Server Monitoring with eBPF and Netflix Hollow"
    p.font.color.rgb = RGBColor(255, 255, 255)  # White text
    p.font.size = Pt(40)
    p.font.bold = True

    # Subtitle
    subtitle_box = slide.shapes.add_textbox(Inches(1), Inches(3), Inches(11.33), Inches(0.5))
    subtitle_text = subtitle_box.text_frame
    subtitle_p = subtitle_text.paragraphs[0]
    subtitle_p.text = "Real-time Performance Monitoring and Visualization"
    subtitle_p.font.color.rgb = INFOBLOX_DARK_BLUE
    subtitle_p.font.size = Pt(28)

    # Author and date
    author_box = slide.shapes.add_textbox(Inches(1), Inches(4.5), Inches(11.33), Inches(1))
    author_text = author_box.text_frame
    author_p = author_text.paragraphs[0]
    author_p.text = "[Your Name]"
    author_p.font.color.rgb = INFOBLOX_GRAY
    author_p.font.size = Pt(20)

    # Add date
    date_p = author_text.add_paragraph()
    date_p.text = "May 10, 2025"
    date_p.font.color.rgb = INFOBLOX_GRAY
    date_p.font.size = Pt(16)

    # Add footer
    add_infoblox_footer(slide)

# Function to create Infoblox styled content slides
def create_content_slide(content_slide):
    slide_layout = infoblox_prs.slide_layouts[1]  # Content slide with title
    slide = infoblox_prs.slides.add_slide(slide_layout)

    # Get the title from original slide
    title_text = ""
    for shape in content_slide.shapes:
        if hasattr(shape, "text") and shape.text.strip():
            title_text = shape.text.strip()
            break

    # Create title with Infoblox styling
    title_shape = slide.shapes.add_shape(
        1, Inches(0), Inches(0), Inches(13.33), Inches(0.85)
    )
    title_shape.fill.solid()
    title_shape.fill.fore_color.rgb = INFOBLOX_BLUE
    title_shape.line.fill.background()

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.1), Inches(12.33), Inches(0.65))
    title_text_frame = title_box.text_frame
    title_p = title_text_frame.paragraphs[0]
    title_p.text = title_text
    title_p.font.color.rgb = RGBColor(255, 255, 255)  # White text
    title_p.font.size = Pt(28)
    title_p.font.bold = True

    # Extract and process content from the original slide
    content_text = []
    for shape in content_slide.shapes:
        if hasattr(shape, "text_frame"):
            for paragraph in shape.text_frame.paragraphs:
                text = paragraph.text.strip()
                if text and text != title_text:
                    content_text.append((text, paragraph.level))

    # Create content area
    content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.0), Inches(12.33), Inches(5.5))
    content_text_frame = content_box.text_frame
    content_text_frame.word_wrap = True

    # Process content and create styled paragraphs
    current_section = None
    first_paragraph = True

    for text, level in content_text:
        # Skip the title which we already processed
        if text == title_text:
            continue

        # Check if this is a section header (bold, ends with :)
        is_section_header = text.endswith(':') or (text.endswith(':') and level == 0)

        if first_paragraph:
            p = content_text_frame.paragraphs[0]
            first_paragraph = False
        else:
            p = content_text_frame.add_paragraph()

        p.text = text
        p.level = level

        if is_section_header:
            p.font.bold = True
            p.font.color.rgb = INFOBLOX_DARK_BLUE
            p.font.size = Pt(24)
            current_section = text
            # Add some space after section headers
            p.space_after = Pt(12)
        else:
            p.font.color.rgb = INFOBLOX_GRAY
            p.font.size = Pt(20)
            p.space_after = Pt(6)

    # Add footer
    add_infoblox_footer(slide)

# Process each slide
for i, slide in enumerate(prs.slides):
    if i == 0:
        # Special handling for title slide
        create_title_slide(slide)
    else:
        # Content slides
        create_content_slide(slide)

# Special handling for code blocks in slide 6 (Demo & Usage)
# Find slide 6 (index 5)
demo_slide = infoblox_prs.slides[5]  # 0-based indexing, so slide 6 is index 5

# Find content text box and extract content
content_box = None
for shape in demo_slide.shapes:
    if hasattr(shape, "text_frame"):
        for paragraph in shape.text_frame.paragraphs:
            if "Running the System:" in paragraph.text:
                content_box = shape
                break
        if content_box:
            break

# If we found the content box with "Running the System:", add code block
if content_box:
    # Add code block with gray background
    code_box = demo_slide.shapes.add_shape(
        1, Inches(1.0), Inches(2.2), Inches(11.33), Inches(1.5)
    )
    code_box.fill.solid()
    code_box.fill.fore_color.rgb = INFOBLOX_LIGHT_GRAY
    code_box.line.color.rgb = INFOBLOX_GRAY

    code_text = demo_slide.shapes.add_textbox(
        Inches(1.2), Inches(2.3), Inches(11.0), Inches(1.3)
    )
    tf = code_text.text_frame

    p1 = tf.paragraphs[0]
    p1.text = "# Start the demo with default settings"
    p1.font.name = "Courier New"
    p1.font.size = Pt(16)
    p1.font.color.rgb = INFOBLOX_DARK_BLUE

    p2 = tf.add_paragraph()
    p2.text = "./run_demo.sh"
    p2.font.name = "Courier New"
    p2.font.size = Pt(16)
    p2.font.color.rgb = INFOBLOX_DARK_BLUE

# Save the styled presentation
output_path = '/home/parallels/Work/Tutorials/Hackathon/2025/Kea_DHCP_Monitoring_Infoblox_Style.pptx'
infoblox_prs.save(output_path)
print(f"Infoblox styled presentation created at: {output_path}")
