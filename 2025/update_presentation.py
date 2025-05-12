#!/usr/bin/env python3
"""
Script to update PowerPoint slides by removing no-fallback option references
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# Create a presentation object
prs = Presentation()

# Set slide layout to widescreen
prs.slide_width = Inches(13.33)
prs.slide_height = Inches(7.5)

# Define colors
TITLE_COLOR = RGBColor(68, 114, 196)  # Blue
SUBTITLE_COLOR = RGBColor(89, 89, 89)  # Dark gray
ACCENT_COLOR = RGBColor(255, 140, 0)   # Orange

# Slide 1: Title Slide
slide_layout = prs.slide_layouts[0]  # Title slide layout
slide = prs.slides.add_slide(slide_layout)

# Title
title = slide.shapes.title
title.text = "Kea DHCP Server Monitoring with eBPF and Netflix Hollow"
title_tf = title.text_frame.paragraphs[0].font
title_tf.size = Pt(40)
title_tf.color.rgb = TITLE_COLOR
title_tf.bold = True

# Subtitle
subtitle = slide.placeholders[1]
subtitle.text = "Real-time Performance Monitoring and Visualization"
subtitle_tf = subtitle.text_frame.paragraphs[0].font
subtitle_tf.size = Pt(28)
subtitle_tf.color.rgb = SUBTITLE_COLOR

# Author and date
author_box = slide.shapes.add_textbox(Inches(0.5), Inches(5.5), Inches(12), Inches(1))
author_text = author_box.text_frame
author_text.paragraphs[0].text = "Your Name"
author_text.paragraphs[0].alignment = PP_ALIGN.LEFT
author_text.add_paragraph().text = "May 10, 2025"

# Slide 2: Overview & Problem
slide_layout = prs.slide_layouts[1]  # Title and Content layout
slide = prs.slides.add_slide(slide_layout)
slide.shapes.title.text = "Overview & Problem"

# Content
content = slide.placeholders[1].text_frame
p = content.paragraphs[0]
p.text = "Challenge:"
p.font.bold = True
p.font.size = Pt(24)
p.font.color.rgb = ACCENT_COLOR

bullet_points = [
    "Need real-time monitoring of Kea DHCP server performance",
    "Must collect metrics without modifying source code",
    "Need efficient storage and visualization of metrics data"
]
for point in bullet_points:
    p = content.add_paragraph()
    p.text = point
    p.level = 1
    p.font.size = Pt(20)

p = content.add_paragraph()
p.text = "\nSolution:"
p.font.bold = True
p.font.size = Pt(24)
p.font.color.rgb = ACCENT_COLOR

solution_points = [
    "eBPF-based monitoring system with Netflix Hollow integration",
    "Real-time dashboard with multiple visualization options",
    "Resilient architecture with fallback mechanisms"
]
for point in solution_points:
    p = content.add_paragraph()
    p.text = point
    p.level = 1
    p.font.size = Pt(20)

# Slide 3: Technology Stack
slide_layout = prs.slide_layouts[1]  # Title and Content layout
slide = prs.slides.add_slide(slide_layout)
slide.shapes.title.text = "Technology Stack"

# Content
content = slide.placeholders[1].text_frame
p = content.paragraphs[0]
p.text = "eBPF Technology:"
p.font.bold = True
p.font.size = Pt(24)
p.font.color.rgb = ACCENT_COLOR

ebpf_points = [
    "Attaches probes to Kea DHCP server functions",
    "Collects metrics with minimal overhead",
    "Monitors: packet processing times, CPU/memory usage, error rates"
]
for point in ebpf_points:
    p = content.add_paragraph()
    p.text = point
    p.level = 1
    p.font.size = Pt(20)

p = content.add_paragraph()
p.text = "\nNetflix Hollow Integration:"
p.font.bold = True
p.font.size = Pt(24)
p.font.color.rgb = ACCENT_COLOR

hollow_points = [
    "Efficient in-memory data storage",
    "Fast data access and versioning",
    "Structured metrics schema"
]
for point in hollow_points:
    p = content.add_paragraph()
    p.text = point
    p.level = 1
    p.font.size = Pt(20)

# Slide 4: System Architecture
slide_layout = prs.slide_layouts[1]  # Title and Content layout
slide = prs.slides.add_slide(slide_layout)
slide.shapes.title.text = "System Architecture"

# Content
content = slide.placeholders[1].text_frame
p = content.paragraphs[0]
p.text = "Components & Data Flow:"
p.font.bold = True
p.font.size = Pt(24)
p.font.color.rgb = ACCENT_COLOR

architecture_points = [
    "eBPF Probes → Raw Metrics → JSON → Hollow Format → Dashboard",
    "Core monitoring: kea_metrics.c/.py",
    "Data processing: ebpf_to_json.py, ebpf_to_hollow.py",
    "Continuous monitoring: ebpf_hollow_monitor.py",
    "Visualization: Real-time dashboard"
]
for point in architecture_points:
    p = content.add_paragraph()
    p.text = point
    p.level = 1
    p.font.size = Pt(20)

p = content.add_paragraph()
p.text = "\nResilient Design:"
p.font.bold = True
p.font.size = Pt(24)
p.font.color.rgb = ACCENT_COLOR

resilient_points = [
    "Local mode when Hollow server is unavailable",
    "Sample metrics fallback when needed",
    "Automatic dashboard recovery"
]
for point in resilient_points:
    p = content.add_paragraph()
    p.text = point
    p.level = 1
    p.font.size = Pt(20)

# Slide 5: Resilient Architecture
slide_layout = prs.slide_layouts[1]  # Title and Content layout
slide = prs.slides.add_slide(slide_layout)
slide.shapes.title.text = "Resilient Architecture"

# Content
content = slide.placeholders[1].text_frame
p = content.paragraphs[0]
p.text = "Fallback Mechanisms:"
p.font.bold = True
p.font.size = Pt(24)
p.font.color.rgb = ACCENT_COLOR

fallback_points = [
    "Local mode when Hollow server is unavailable",
    "Sample metrics when collection fails",
    "Automatic dashboard recovery"
]
for point in fallback_points:
    p = content.add_paragraph()
    p.text = point
    p.level = 1
    p.font.size = Pt(20)

p = content.add_paragraph()
p.text = "\nEnhanced Error Handling:"
p.font.bold = True
p.font.size = Pt(24)
p.font.color.rgb = ACCENT_COLOR

error_handling_points = [
    "Clear logging of system state",
    "Graceful degradation when components fail",
    "Fixed working directory issue in metrics collection"
]
for point in error_handling_points:
    p = content.add_paragraph()
    p.text = point
    p.level = 1
    p.font.size = Pt(20)

# Slide 6: Demo & Usage
slide_layout = prs.slide_layouts[1]  # Title and Content layout
slide = prs.slides.add_slide(slide_layout)
slide.shapes.title.text = "Demo & Usage"

# Content
content = slide.placeholders[1].text_frame
p = content.paragraphs[0]
p.text = "Running the System:"
p.font.bold = True
p.font.size = Pt(24)
p.font.color.rgb = ACCENT_COLOR

# Add code block
code_textbox = slide.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(11.5), Inches(1))
code_tf = code_textbox.text_frame
code_p = code_tf.paragraphs[0]
code_p.text = "# Start the demo with default settings"
code_p.font.name = "Courier New"
code_p.font.size = Pt(18)

code_p = code_tf.add_paragraph()
code_p.text = "./run_demo.sh"
code_p.font.name = "Courier New"
code_p.font.size = Pt(18)

# Add key features section
features_textbox = slide.shapes.add_textbox(Inches(0.8), Inches(3.5), Inches(11.5), Inches(2))
features_tf = features_textbox.text_frame
p = features_tf.paragraphs[0]
p.text = "Key Features to Observe:"
p.font.bold = True
p.font.size = Pt(24)
p.font.color.rgb = ACCENT_COLOR

features_points = [
    "Real-time metrics at http://localhost:8080/",
    "Automatic fallback mechanisms in action",
    "Graceful handling of collection failures"
]
for point in features_points:
    p = features_tf.add_paragraph()
    p.text = point
    p.font.size = Pt(20)

# Slide 7: Real-World Benefits
slide_layout = prs.slide_layouts[1]  # Title and Content layout
slide = prs.slides.add_slide(slide_layout)
slide.shapes.title.text = "Real-World Benefits"

# Content
content = slide.placeholders[1].text_frame
p = content.paragraphs[0]
p.text = "Operational Advantages:"
p.font.bold = True
p.font.size = Pt(24)
p.font.color.rgb = ACCENT_COLOR

benefits_points = [
    "Early detection of performance issues",
    "Non-invasive monitoring (no code changes)",
    "Production-ready monitoring solution",
    "Clear visibility into system performance"
]
for point in benefits_points:
    p = content.add_paragraph()
    p.text = point
    p.level = 1
    p.font.size = Pt(20)

p = content.add_paragraph()
p.text = "\nUse Cases:"
p.font.bold = True
p.font.size = Pt(24)
p.font.color.rgb = ACCENT_COLOR

use_cases = [
    "DevOps continuous monitoring",
    "Performance troubleshooting",
    "Capacity planning",
    "Production DHCP server monitoring"
]
for case in use_cases:
    p = content.add_paragraph()
    p.text = case
    p.level = 1
    p.font.size = Pt(20)

# Slide 8: Conclusion & Next Steps
slide_layout = prs.slide_layouts[1]  # Title and Content layout
slide = prs.slides.add_slide(slide_layout)
slide.shapes.title.text = "Conclusion & Next Steps"

# Content
content = slide.placeholders[1].text_frame
p = content.paragraphs[0]
p.text = "Key Takeaways:"
p.font.bold = True
p.font.size = Pt(24)
p.font.color.rgb = ACCENT_COLOR

takeaways = [
    "eBPF + Netflix Hollow provides powerful, efficient monitoring",
    "Resilient architecture ensures reliable operation",
    "Easy to deploy and use in various environments"
]
for point in takeaways:
    p = content.add_paragraph()
    p.text = point
    p.level = 1
    p.font.size = Pt(20)

p = content.add_paragraph()
p.text = "\nFuture Enhancements:"
p.font.bold = True
p.font.size = Pt(24)
p.font.color.rgb = ACCENT_COLOR

enhancements = [
    "Additional eBPF probe points",
    "Integration with alerting systems",
    "Enhanced visualization options"
]
for point in enhancements:
    p = content.add_paragraph()
    p.text = point
    p.level = 1
    p.font.size = Pt(20)

# Thank you slide with contact information
p = content.add_paragraph()
p.text = "\nThank you! Questions?"
p.alignment = PP_ALIGN.CENTER
p.font.bold = True
p.font.size = Pt(28)
p.font.color.rgb = TITLE_COLOR

# Save the presentation
prs.save('/home/parallels/Work/Tutorials/Hackathon/2025/Kea_DHCP_Monitoring_Presentation.pptx')
print("Presentation updated successfully!")
