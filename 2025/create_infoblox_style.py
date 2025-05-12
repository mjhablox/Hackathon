#!/usr/bin/env python3

import traceback
import sys

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

    # Define Infoblox brand colors
    # Primary colors
    INFOBLOX_BLUE = RGBColor(0, 125, 186)       # #007DBA - Primary blue
    INFOBLOX_DARK_BLUE = RGBColor(0, 63, 94)    # #003F5E - Dark blue
    INFOBLOX_LIGHT_BLUE = RGBColor(0, 176, 240) # #00B0F0 - Light blue
    INFOBLOX_GRAY = RGBColor(88, 88, 91)        # #58585B - Dark gray
    INFOBLOX_LIGHT_GRAY = RGBColor(236, 236, 236) # #ECECEC - Light gray background

    # Load the existing presentation
    presentation_path = '/home/parallels/Work/Tutorials/Hackathon/2025/Kea_DHCP_Monitoring_Presentation.pptx'
    print(f"Loading presentation from: {presentation_path}")
    prs = Presentation(presentation_path)
    print(f"Successfully loaded presentation with {len(prs.slides)} slides")

    # Create a new presentation with Infoblox styling
    infoblox_prs = Presentation()
    infoblox_prs.slide_width = Inches(13.33)
    infoblox_prs.slide_height = Inches(7.5)
    print("Created new presentation for Infoblox styling")

    # Function to add an Infoblox footer to a slide
    def add_infoblox_footer(slide):
        # Add a thin line above the footer
        line_shape = slide.shapes.add_shape(
            1, Inches(0.5), Inches(6.9), Inches(12.33), Inches(0.02)
        )
        line_shape.fill.solid()
        line_shape.fill.fore_color.rgb = INFOBLOX_BLUE
        line_shape.line.fill.background()

        # Add footer text
        footer = slide.shapes.add_textbox(Inches(0.5), Inches(7.0), Inches(8), Inches(0.3))
        footer_text = footer.text_frame
        footer_text.paragraphs[0].text = "Kea DHCP Server Monitoring | May 10, 2025"
        footer_text.paragraphs[0].font.size = Pt(10)
        footer_text.paragraphs[0].font.color.rgb = INFOBLOX_GRAY

        # Add logo placeholder (right side of footer)
        logo_placeholder = slide.shapes.add_textbox(Inches(11), Inches(7.0), Inches(1.8), Inches(0.3))
        logo_text = logo_placeholder.text_frame
        logo_text.paragraphs[0].text = "INFOBLOX"
        logo_text.paragraphs[0].alignment = PP_ALIGN.RIGHT
        logo_text.paragraphs[0].font.bold = True
        logo_text.paragraphs[0].font.size = Pt(12)
        logo_text.paragraphs[0].font.color.rgb = INFOBLOX_BLUE

    # Create title slide (first slide)
    slide_layout = infoblox_prs.slide_layouts[0]  # Title Slide layout
    slide = infoblox_prs.slides.add_slide(slide_layout)
    print("Adding title slide")

    # Add blue background shape at the top
    top_bg = slide.shapes.add_shape(
        1, Inches(0), Inches(0), Inches(13.33), Inches(2.5)
    )
    top_bg.fill.solid()
    top_bg.fill.fore_color.rgb = INFOBLOX_BLUE
    top_bg.line.fill.background()

    # Title text
    title_box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(11.33), Inches(1))
    title_text = title_box.text_frame
    title_text.word_wrap = True

    p = title_text.paragraphs[0]
    p.text = "Kea DHCP Server Monitoring with eBPF and Netflix Hollow"
    p.font.color.rgb = RGBColor(255, 255, 255)  # White text
    p.font.size = Pt(40)
    p.font.bold = True

    # Subtitle
    subtitle_box = slide.shapes.add_textbox(Inches(1), Inches(3), Inches(11.33), Inches(0.5))
    subtitle_text = subtitle_box.text_frame
    subtitle_p = subtitle_text.paragraphs[0]
    subtitle_p.text = "Real-time Performance Monitoring and Visualization"
    subtitle_p.font.color.rgb = INFOBLOX_DARK_BLUE
    subtitle_p.font.size = Pt(28)

    # Author and date
    author_box = slide.shapes.add_textbox(Inches(1), Inches(4.5), Inches(11.33), Inches(1))
    author_text = author_box.text_frame
    author_p = author_text.paragraphs[0]
    author_p.text = "[Your Name]"
    author_p.font.color.rgb = INFOBLOX_GRAY
    author_p.font.size = Pt(20)

    # Add date
    date_p = author_text.add_paragraph()
    date_p.text = "May 10, 2025"
    date_p.font.color.rgb = INFOBLOX_GRAY
    date_p.font.size = Pt(16)

    # Add footer
    add_infoblox_footer(slide)

    # Create content slides
    slide_titles = [
        "Overview & Problem",
        "Technology Stack",
        "System Architecture",
        "Resilient Architecture",
        "Demo & Usage",
        "Real-World Benefits",
        "Conclusion & Next Steps"
    ]

    # Content for each slide (simplified)
    slide_content = [
        [
            ["Challenge:", 0, True],
            ["Need real-time monitoring of Kea DHCP server performance", 1, False],
            ["Must collect metrics without modifying source code", 1, False],
            ["Need efficient storage and visualization of metrics data", 1, False],
            ["Solution:", 0, True],
            ["eBPF-based monitoring system with Netflix Hollow integration", 1, False],
            ["Real-time dashboard with multiple visualization options", 1, False],
            ["Resilient architecture with fallback mechanisms", 1, False]
        ],
        [
            ["eBPF Technology:", 0, True],
            ["Attaches probes to Kea DHCP server functions", 1, False],
            ["Collects metrics with minimal overhead", 1, False],
            ["Monitors: packet processing times, CPU/memory usage, error rates", 1, False],
            ["Netflix Hollow Integration:", 0, True],
            ["Efficient in-memory data storage", 1, False],
            ["Fast data access and versioning", 1, False],
            ["Structured metrics schema", 1, False]
        ],
        [
            ["Components & Data Flow:", 0, True],
            ["eBPF Probes → Raw Metrics → JSON → Hollow Format → Dashboard", 1, False],
            ["Core monitoring: kea_metrics.c/.py", 1, False],
            ["Data processing: ebpf_to_json.py, ebpf_to_hollow.py", 1, False],
            ["Continuous monitoring: ebpf_hollow_monitor.py", 1, False],
            ["Visualization: Real-time dashboard", 1, False],
            ["Resilient Design:", 0, True],
            ["Local mode when Hollow server is unavailable", 1, False],
            ["Sample metrics fallback (when needed)", 1, False],
            ["Automatic dashboard recovery", 1, False]
        ],
        [
            ["Fallback Mechanisms:", 0, True],
            ["Local mode when Hollow server is unavailable", 1, False],
            ["Sample metrics when collection fails", 1, False],
            ["Automatic dashboard recovery", 1, False],
            ["Enhanced Error Handling:", 0, True],
            ["Clear logging of system state", 1, False],
            ["Graceful degradation when components fail", 1, False],
            ["Fixed working directory issue in metrics collection", 1, False]
        ],
        [
            ["Running the System:", 0, True],
            ["Start the demo with default settings", 1, False],
            ["./run_demo.sh", 1, False],
            ["Key Features to Observe:", 0, True],
            ["Real-time metrics at http://localhost:8080/", 1, False],
            ["Automatic fallback mechanisms in action", 1, False],
            ["Graceful handling of collection failures", 1, False]
        ],
        [
            ["Operational Advantages:", 0, True],
            ["Early detection of performance issues", 1, False],
            ["Non-invasive monitoring (no code changes)", 1, False],
            ["Production-ready monitoring solution", 1, False],
            ["Clear visibility into system performance", 1, False],
            ["Use Cases:", 0, True],
            ["DevOps continuous monitoring", 1, False],
            ["Performance troubleshooting", 1, False],
            ["Capacity planning", 1, False],
            ["Production DHCP server monitoring", 1, False]
        ],
        [
            ["Key Takeaways:", 0, True],
            ["eBPF + Netflix Hollow provides powerful, efficient monitoring", 1, False],
            ["Resilient architecture ensures reliable operation", 1, False],
            ["Easy to deploy and use in various environments", 1, False],
            ["Future Enhancements:", 0, True],
            ["Additional eBPF probe points", 1, False],
            ["Integration with alerting systems", 1, False],
            ["Enhanced visualization options", 1, False],
            ["Thank you! Questions?", 0, True]
        ]
    ]

    for i, (title, content_items) in enumerate(zip(slide_titles, slide_content)):
        print(f"Creating slide {i+2}: {title}")

        # Create the slide
        slide_layout = infoblox_prs.slide_layouts[1]  # Content slide layout
        slide = infoblox_prs.slides.add_slide(slide_layout)

        # Create title with Infoblox styling
        title_shape = slide.shapes.add_shape(
            1, Inches(0), Inches(0), Inches(13.33), Inches(0.85)
        )
        title_shape.fill.solid()
        title_shape.fill.fore_color.rgb = INFOBLOX_BLUE
        title_shape.line.fill.background()

        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.1), Inches(12.33), Inches(0.65))
        title_text_frame = title_box.text_frame
        title_p = title_text_frame.paragraphs[0]
        title_p.text = title
        title_p.font.color.rgb = RGBColor(255, 255, 255)  # White text
        title_p.font.size = Pt(28)
        title_p.font.bold = True

        # Create content area
        content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.0), Inches(12.33), Inches(5.5))
        content_tf = content_box.text_frame
        content_tf.word_wrap = True

        # Add content items
        first_para = True
        current_header = None

        for text, level, is_header in content_items:
            # Skip adding the content to the demo slide - we'll handle that specially
            if title == "Demo & Usage" and text == "./run_demo.sh":
                continue

            if first_para:
                p = content_tf.paragraphs[0]
                first_para = False
            else:
                p = content_tf.add_paragraph()

            p.text = text
            p.level = level

            if is_header:
                p.font.bold = True
                p.font.color.rgb = INFOBLOX_DARK_BLUE
                p.font.size = Pt(24)
                current_header = text
                p.space_after = Pt(12)  # Add spacing after headers
            else:
                p.font.color.rgb = INFOBLOX_GRAY
                p.font.size = Pt(20)
                p.space_after = Pt(6)

        # Add footer
        add_infoblox_footer(slide)

        # Special handling for Demo & Usage slide (slide 6)
        if title == "Demo & Usage":
            # Add code block with gray background
            code_box = slide.shapes.add_shape(
                1, Inches(0.8), Inches(2.5), Inches(11.5), Inches(1.0)
            )
            code_box.fill.solid()
            code_box.fill.fore_color.rgb = INFOBLOX_LIGHT_GRAY
            code_box.line.color.rgb = INFOBLOX_GRAY

            # Add the code text
            code_text = slide.shapes.add_textbox(
                Inches(1.0), Inches(2.6), Inches(11.0), Inches(0.8)
            )
            code_tf = code_text.text_frame

            code_p1 = code_tf.paragraphs[0]
            code_p1.text = "# Start the demo with default settings"
            code_p1.font.name = "Courier New"
            code_p1.font.size = Pt(16)
            code_p1.font.color.rgb = INFOBLOX_DARK_BLUE

            code_p2 = code_tf.add_paragraph()
            code_p2.text = "./run_demo.sh"
            code_p2.font.name = "Courier New"
            code_p2.font.size = Pt(16)
            code_p2.font.color.rgb = INFOBLOX_DARK_BLUE

    # Save the presentation
    output_path = '/home/parallels/Work/Tutorials/Hackathon/2025/Kea_DHCP_Monitoring_Infoblox_Style.pptx'
    print(f"Saving Infoblox styled presentation to: {output_path}")
    infoblox_prs.save(output_path)
    print(f"Infoblox styled presentation created successfully!")

except Exception as e:
    print(f"Error occurred: {e}")
    traceback.print_exc()
    sys.exit(1)
